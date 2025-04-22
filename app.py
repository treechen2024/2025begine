from flask import Flask, render_template, request, jsonify, send_file, url_for
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import uuid
import re
from google.generativeai import GenerativeModel, configure
from dotenv import load_dotenv
import os

# 載入環境變數
load_dotenv()

app = Flask(__name__)
app.config['JSON_AS_ASCII'] = False
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB 上傳限制

# 確保上傳和輸出目錄存在
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    if 'txtFile' not in request.files:
        return jsonify({'success': False, 'error': '未找到檔案'})
    
    txt_file = request.files['txtFile']
    if txt_file.filename == '':
        return jsonify({'success': False, 'error': '未選擇檔案'})
    
    # 獲取顏色值
    font_color = request.form.get('fontColor', '#000000')
    bg_color = request.form.get('bgColor', '#FFFFFF')
    
    # 轉換十六進制顏色為RGB元組
    font_rgb = hex_to_rgb(font_color)
    bg_rgb = hex_to_rgb(bg_color)
    
    # 生成唯一檔案名
    txt_filename = str(uuid.uuid4()) + '.txt'
    txt_path = os.path.join(app.config['UPLOAD_FOLDER'], txt_filename)
    
    # 保存上傳的TXT檔案
    txt_file.save(txt_path)
    
    # 生成PPT輸出路徑
    output_filename = str(uuid.uuid4()) + '.pptx'
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    
    try:
        # 創建PPT
        create_pptx(txt_path, output_path, font_rgb, bg_rgb)
        
        # 生成下載URL
        download_url = url_for('download_file', filename=output_filename)
        
        return jsonify({
            'success': True,
            'download_url': download_url
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})
    finally:
        # 清理上傳的TXT檔案
        if os.path.exists(txt_path):
            os.remove(txt_path)

@app.route('/download/<filename>')
def download_file(filename):
    return send_file(os.path.join(app.config['OUTPUT_FOLDER'], filename),
                     as_attachment=True,
                     download_name='presentation.pptx')

def hex_to_rgb(hex_color):
    """將十六進制顏色代碼轉換為RGB元組"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_pptx(txt_path, output_path, font_color, bg_color):
    """創建PPT檔案"""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    with open(txt_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # 以兩個以上空白行為分頁依據
    blocks = [block.strip() for block in re.split(r'\n\s*\n', content) if block.strip()]

    for block in blocks:
        slide = prs.slides.add_slide(blank_slide_layout)

        # 背景色
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

        # 新增文字方塊
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
        tf = textbox.text_frame
        tf.clear()

        # 使用 /N 作為內頁換行
        lines = block.split('/N')
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line.strip()
            run = p.runs[0]
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(*font_color)

    prs.save(output_path)

@app.route('/generate', methods=['POST'])
def generate():
    try:
        prompt = request.json.get('prompt', '')
        if not prompt:
            return jsonify({'error': '請輸入提示語！'}), 400
        
        # 添加超時處理
        response = model.generate_content(prompt, timeout=30)
        if not response or not response.text:
            return jsonify({'error': '生成內容失敗，請重試'}), 500
            
        return jsonify({'result': response.text})
    except Exception as e:
        app.logger.error(f"生成內容時發生錯誤: {str(e)}")
        return jsonify({'error': f'發生錯誤：{str(e)}'}), 500

@app.errorhandler(404)
def not_found(error):
    return jsonify({'error': '找不到該頁面'}), 404

@app.errorhandler(500)
def server_error(error):
    return jsonify({'error': '伺服器錯誤'}), 500

if __name__ == '__main__':
    port = int(os.getenv('PORT', 5000))
    app.run(host='0.0.0.0', port=port)